"""PPTX parser adapter using python-pptx."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from adx.models.document import FileType
from adx.parsers.base import (
    ParserAdapter,
    ParserCapabilities,
    ParserResult,
    ParserWarning,
)

logger = logging.getLogger(__name__)


class PptxAdapter(ParserAdapter):
    """Parser adapter for PPTX files using python-pptx."""

    def name(self) -> str:
        return "python-pptx"

    def version(self) -> str:
        try:
            import pptx

            return getattr(pptx, "__version__", "unknown")
        except ImportError:
            return "not installed"

    def capabilities(self) -> ParserCapabilities:
        return ParserCapabilities(
            supported_types=[FileType.PPTX],
            extracts_text=True,
            extracts_tables=True,
            extracts_layout=False,
            extracts_formulas=False,
            extracts_images=False,
            supports_ocr=False,
        )

    def parse(self, file_path: Path, file_type: FileType) -> ParserResult:
        result = ParserResult(
            parser_name=self.name(),
            parser_version=self.version(),
            file_type=file_type,
        )

        try:
            from pptx import Presentation
        except ImportError:
            result.errors.append(
                "python-pptx is not installed. Run: pip install python-pptx"
            )
            return result

        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            result.errors.append(f"Failed to open PPTX: {e}")
            return result

        for slide_num, slide in enumerate(prs.slides, 1):
            text_blocks: list[dict[str, Any]] = []
            tables: list[dict[str, Any]] = []
            block_index = 0

            title_text = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()

            for shape in slide.shapes:
                if shape.has_text_frame:
                    is_title_shape = (
                        slide.shapes.title is not None
                        and shape.shape_id == slide.shapes.title.shape_id
                    )
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue

                        if is_title_shape:
                            block_type = "heading"
                            font_size = 24.0
                        else:
                            block_type = "paragraph"
                            font_size = 12.0

                        text_blocks.append({
                            "text": text,
                            "block_type": block_type,
                            "page_number": slide_num,
                            "font_size": font_size,
                            "confidence": 1.0,
                            "reading_order_index": block_index,
                        })
                        block_index += 1

                if shape.has_table:
                    table = shape.table
                    cells: list[dict[str, Any]] = []
                    row_count = len(table.rows)
                    col_count = len(table.columns)

                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cells.append({
                                "row_index": row_idx,
                                "column_index": col_idx,
                                "value": cell.text.strip(),
                                "data_type": "string",
                            })

                    tables.append({
                        "page_number": slide_num,
                        "title": f"Table on Slide {slide_num}",
                        "row_count": row_count,
                        "column_count": col_count,
                        "cells": cells,
                        "confidence": 1.0,
                    })

            # Speaker notes as footnote blocks
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text_blocks.append({
                        "text": notes_text,
                        "block_type": "footnote",
                        "page_number": slide_num,
                        "font_size": 10.0,
                        "confidence": 1.0,
                        "reading_order_index": block_index,
                    })
                    block_index += 1

            page_data: dict[str, Any] = {
                "page_number": slide_num,
                "width": prs.slide_width.pt if prs.slide_width else 960,
                "height": prs.slide_height.pt if prs.slide_height else 540,
                "text_blocks": text_blocks,
                "tables": tables,
            }

            result.pages.append(page_data)
            result.text_blocks.extend(text_blocks)
            result.tables.extend(tables)

        result.page_count = len(prs.slides)

        # Extract metadata
        try:
            props = prs.core_properties
            result.metadata = {
                "title": props.title or "",
                "author": props.author or "",
                "subject": props.subject or "",
                "created": str(props.created) if props.created else "",
                "modified": str(props.modified) if props.modified else "",
                "slide_count": len(prs.slides),
            }
        except Exception:
            result.metadata = {"slide_count": len(prs.slides)}

        if not result.text_blocks and not result.tables:
            result.warnings.append(
                ParserWarning(
                    code="EMPTY_DOCUMENT",
                    message="PPTX contains no text or tables.",
                )
            )

        return result
