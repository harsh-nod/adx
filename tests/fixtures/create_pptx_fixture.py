"""Create a sample PPTX fixture for testing."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def create_sample_pptx():
    """Create a PPTX with a title slide, content slide with bullets, and a table slide with notes."""
    prs = Presentation()

    prs.core_properties.title = "Test Presentation"
    prs.core_properties.author = "Test Author"
    prs.core_properties.subject = "Testing"

    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Test Presentation Title"
    slide1.placeholders[1].text = "Subtitle text for testing"

    # Slide 2: Content slide with bullets
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Content Slide"
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "First bullet point"
    p = tf.add_paragraph()
    p.text = "Second bullet point"
    p = tf.add_paragraph()
    p.text = "Third bullet point"

    # Slide 3: Table slide with speaker notes
    slide_layout = prs.slide_layouts[5]  # Blank
    slide3 = prs.slides.add_slide(slide_layout)

    # Add a title textbox
    from pptx.util import Emu
    txBox = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    txBox.text_frame.text = "Table Slide"

    # Add a table
    rows, cols = 4, 3
    table_shape = slide3.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3))
    table = table_shape.table

    headers = ["Item", "Quantity", "Price"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h

    data = [
        ["Widget A", "10", "25.00"],
        ["Widget B", "5", "50.00"],
        ["Total", "", "500.00"],
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            table.cell(row_idx + 1, col_idx).text = val

    # Speaker notes
    notes_slide = slide3.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes for the table slide."

    output_dir = Path(__file__).parent / "pptx"
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "sample.pptx"
    prs.save(str(output_path))
    print(f"Created {output_path}")


if __name__ == "__main__":
    create_sample_pptx()
